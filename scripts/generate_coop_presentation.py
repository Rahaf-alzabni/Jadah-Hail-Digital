"""Generate COOP Open Day presentation from University of Hail template."""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = Path(r'c:\Users\a\OneDrive\Documents\9 COOP_Presentation_Template.pptx')
OUTPUT = ROOT / 'docs' / 'COOP_Presentation_Jadah_Hail.pptx'
DIAGRAMS = ROOT / 'docs' / 'diagrams'
SCREENSHOTS = ROOT / 'docs' / 'screenshots'

STUDENT_NAME = 'Rahaf Bader Alzabni'
STUDENT_ID = 's202005878'
SUPERVISOR = '[Supervisor Name]'  # Edit before presenting
COOP_TITLE = (
    'Jadah Hail Digital: Smart Tourism Guide\n'
    'Integrated with Tawakkalna Municipality Services'
)
EMPLOYER = 'Hail Municipality\nDigital Transformation Agency'


def set_text(shape, text, font_size=18, bold=False):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold


def set_bullets(shape, items, font_size=16):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(font_size)


def find_content_placeholder(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Content Placeholder'):
            return shape
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name not in ('Title 1',):
            if 'شعار' not in shape.text_frame.text and 'جهة' not in shape.text_frame.text:
                return shape
    return None


def find_title(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and 'Title' in shape.name:
            return shape
    return None


def update_employer_logo(slide, text):
    for shape in slide.shapes:
        if shape.has_text_frame and ('شعار' in shape.text_frame.text or shape.name == 'Rectangle 6'):
            set_text(shape, text, font_size=14, bold=True)
            return


def add_slide(prs, layout_idx, title, bullets=None, image_path=None, image_caption=None):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    update_employer_logo(slide, EMPLOYER.replace('\n', ' | '))

    title_shape = find_title(slide)
    if title_shape:
        set_text(title_shape, title, font_size=28, bold=True)

    content = find_content_placeholder(slide)
    if bullets and content:
        set_bullets(content, bullets, font_size=15)

    if image_path and image_path.exists():
        # Picture with caption layout or blank area
        left = Inches(0.8)
        top = Inches(1.6)
        width = Inches(8.5)
        slide.shapes.add_picture(str(image_path), left, top, width=width)
        if image_caption and content:
            set_text(content, image_caption, font_size=12)

    return slide


def build_presentation():
    shutil.copy(TEMPLATE, OUTPUT)
    prs = Presentation(str(OUTPUT))

    # ── Slide 1: Title ──────────────────────────────────────────────────────
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt == 'Student Name' and shape.name == 'TextBox 7':
            set_text(shape, STUDENT_NAME, font_size=22, bold=True)
        elif txt == 'COOP Title':
            set_text(shape, COOP_TITLE, font_size=18, bold=True)
        elif txt == 'Student ID' and shape.name == 'TextBox 10':
            set_text(shape, STUDENT_ID, font_size=18)
        elif txt == 'Supervisor Name':
            set_text(shape, SUPERVISOR, font_size=18)
        elif txt == 'COOP OPEN DAY DISCUSSIONS':
            set_text(shape, 'COOP OPEN DAY PRESENTATION', font_size=20, bold=True)
    update_employer_logo(s1, EMPLOYER)

    # ── Slide 2: Introduction ───────────────────────────────────────────────
    s2 = prs.slides[1]
    update_employer_logo(s2, EMPLOYER.replace('\n', ' | '))
    content = find_content_placeholder(s2)
    set_bullets(content, [
        'Training at: Hail Municipality — Technology Department',
        'Unit: Digital Transformation Agency',
        'National context: Vision 2030 digital government services',
        'Service channel: Tawakkalna municipality digital services',
        'Training stages:',
        '  • Orientation & qualification (IT governance, service standards)',
        '  • Technical work (data collection, development, testing, reporting)',
        'Main project: Jadah Hail Digital — smart tourism guide for Hail region',
    ], font_size=15)

    # ── Slide 3: Tasks During COOP Training ───────────────────────────────
    s3 = prs.slides[2]
    update_employer_logo(s3, EMPLOYER.replace('\n', ' | '))
    content = find_content_placeholder(s3)
    set_bullets(content, [
        'Tasks Completed:',
        '  • Collected and verified tourism data for Hail landmarks',
        '  • Documented bilingual content (Arabic / English)',
        '  • Developed REST API using Django REST Framework',
        '  • Built visitor portal using React + TypeScript + Vite',
        '  • Integrated interactive map (Leaflet + OpenStreetMap)',
        '  • Implemented reviews, favorites, and tourism assistant',
        '  • Prepared technical reports and COOP documentation',
        '',
        'Issues & Solutions:',
        '  • Issue: Scattered tourism data → Solution: Centralized database & API',
        '  • Issue: Bilingual UI complexity → Solution: RTL/LTR localization layer',
        '  • Issue: Map accuracy → Solution: Verified GPS coordinates per landmark',
    ], font_size=14)

    # ── Slide 4: Experiences from COOP Training ─────────────────────────────
    s4 = prs.slides[3]
    update_employer_logo(s4, EMPLOYER.replace('\n', ' | '))
    content = find_content_placeholder(s4)
    set_bullets(content, [
        'Professional skills gained:',
        '  • Full-stack web development (Django + React)',
        '  • REST API design and session-based authentication',
        '  • Data collection, organization, and technical reporting',
        '  • Working within government digital service standards',
        '',
        'Soft skills & workplace experience:',
        '  • Organization and structured reporting in a municipal environment',
        '  • Innovation within institutional requirements',
        '  • Attention to data accuracy for public-facing services',
        '  • Understanding Tawakkalna as a national service delivery channel',
    ], font_size=14)

    # ── Slide 5: COOP Project — Results ─────────────────────────────────────
    s5 = prs.slides[4]
    update_employer_logo(s5, EMPLOYER.replace('\n', ' | '))
    title = find_title(s5)
    if title:
        set_text(title, 'COOP Project: Jadah Hail Digital — Results', font_size=26, bold=True)
    content = find_content_placeholder(s5)
    set_bullets(content, [
        'Delivered: Bilingual smart tourism platform for Hail Municipality',
        'Content: 6 landmarks, 3 events, 3 routes, tourism assistant',
        'Features: Interactive map, reviews, favorites, admin dashboard',
        'Architecture: Three-tier (React → Django REST API → SQLite)',
        'Testing: 19 automated unit tests — all passed',
        'Integration: Designed as municipality service within Tawakkalna ecosystem',
        'Repository: github.com/moneerafahaid-collab/WATN-2-',
    ], font_size=14)

    # ── Additional slides ───────────────────────────────────────────────────
    add_slide(prs, 1, 'System Architecture', [
        'Presentation Layer: React + TypeScript + Tailwind CSS',
        'Application Layer: Django 6 + Django REST Framework',
        'Data Layer: SQLite + media files for images',
        'External: OpenStreetMap tiles for interactive map',
        'Key API endpoints: /places/, /events/, /routes/, /reviews/, /auth/',
    ])
    arch_slide = prs.slides[-1]
    arch_slide.shapes.add_picture(
        str(DIAGRAMS / 'architecture_diagram.png'),
        Inches(5.2), Inches(1.5), width=Inches(4.5),
    )

    add_slide(prs, 1, 'Use Case Overview', [
        'Actors: Visitor, Registered User, Municipality Admin',
        'Visitor: browse places, map, events, routes, assistant',
        'Registered User: login, favorites, submit reviews',
        'Admin: dashboard, content management, analytics overview',
    ])
    prs.slides[-1].shapes.add_picture(
        str(DIAGRAMS / 'use_case_diagram.png'),
        Inches(4.8), Inches(1.4), width=Inches(4.8),
    )

    # Demo screenshots
    demo_slides = [
        ('01_home.png', 'Visitor Portal — Home Page'),
        ('03_map.png', 'Interactive Hail Tourism Map'),
        ('06_place_details.png', 'Landmark Details & Reviews'),
        ('07_admin_dashboard.png', 'Municipality Admin Dashboard'),
    ]
    for img, caption in demo_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        update_employer_logo(slide, EMPLOYER.replace('\n', ' | '))
        title = find_title(slide)
        if title:
            set_text(title, caption, font_size=24, bold=True)
        img_path = SCREENSHOTS / img
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(0.6), Inches(1.3), width=Inches(8.8))

    add_slide(prs, 1, 'Testing & Validation', [
        'Manual testing: all portal pages and admin interface — Passed',
        'API testing: places, events, routes, auth — Passed',
        'Localization: Arabic RTL / English LTR — Passed',
        'Automated unit tests: 19 tests (tourism, accounts, events)',
        'Command: python manage.py test tourism accounts events',
        'Result: OK — 19/19 tests passed',
    ])

    add_slide(prs, 1, 'Conclusions & Recommendations', [
        'Successfully delivered Jadah Hail Digital as a COOP training outcome',
        'Aligned municipal tourism promotion with digital transformation goals',
        'Recommendations:',
        '  • Complete deep Tawakkalna service integration',
        '  • Deploy on production server with PostgreSQL',
        '  • Connect admin dashboard CRUD to live API',
        '  • Expand automated testing to frontend components',
        '',
        'Thank you — Questions?',
    ])

    # Thank you slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt == 'Student Name' and 'TextBox' in shape.name:
                set_text(shape, STUDENT_NAME, font_size=24, bold=True)
            elif txt == 'COOP Title':
                set_text(shape, 'Thank You\nQuestions & Discussion', font_size=22, bold=True)
            elif txt == 'Student ID' and 'TextBox' in shape.name:
                set_text(shape, f'{STUDENT_ID} | {STUDENT_NAME}', font_size=16)
    update_employer_logo(slide, EMPLOYER)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == '__main__':
    path = build_presentation()
    print(f'Created: {path}')
    print(f'Slides: {len(Presentation(str(path)).slides)}')
    print('Edit SUPERVISOR name in scripts/generate_coop_presentation.py if needed.')
